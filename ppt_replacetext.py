'''
Created on Oct 1, 2017
@author: fufufanatic
'''

from pptx import Presentation

testString = 'find text'
replaceString = 'replace text'
ppt = Presentation('enter path to powerpoint presentation')

def replaceText():
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(testString) != -1):
                    #for debugging purposes - prints index of found text
                    print(shape.text.find(testString))
                    shape.text = shape.text.replace(testString, replaceString)
                    
    ppt.save('enter same path to overwrite file or new path for newly created file')
                            
if __name__ == "__main__":
    replaceText()